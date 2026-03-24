"""RAG layer: load corpus, chunk, embed, and retrieve."""

import os
import shutil
import tempfile
import zipfile

import chromadb
from sentence_transformers import SentenceTransformer

CORPUS_DIR = os.environ.get("CORPUS_DIR", "corpus")
CHROMA_DIR = os.environ.get("CHROMA_DIR", "chroma_data")
CHUNK_SIZE = 500   # approximate token count (words used as proxy)
CHUNK_OVERLAP = 50
TOP_K = 3

_model: SentenceTransformer | None = None
_collection: chromadb.Collection | None = None
_client: chromadb.ClientAPI | None = None

SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".pptx", ".ppt"}


def _get_model() -> SentenceTransformer:
    global _model
    if _model is None:
        _model = SentenceTransformer("all-MiniLM-L6-v2")
    return _model


def _get_client() -> chromadb.ClientAPI:
    global _client
    if _client is None:
        _client = chromadb.Client(chromadb.config.Settings(
            persist_directory=CHROMA_DIR,
            anonymized_telemetry=False,
            is_persistent=True,
        ))
    return _client


def _approximate_token_split(text: str, size: int, overlap: int) -> list[str]:
    """Split text into chunks of approximately `size` words with `overlap`."""
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + size
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start = end - overlap
    return chunks


def _read_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _read_pdf(path: str) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)
    except Exception:
        return ""


def _read_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
        return "\n".join(texts)
    except Exception:
        return ""


def _read_file(path: str) -> str:
    """Read a file based on its extension."""
    lower = path.lower()
    if lower.endswith(".txt"):
        return _read_txt(path)
    elif lower.endswith(".pdf"):
        return _read_pdf(path)
    elif lower.endswith((".pptx", ".ppt")):
        return _read_pptx(path)
    return ""


def _extract_zip(zip_bytes: bytes) -> list[tuple[str, bytes]]:
    """Extract supported files from a ZIP archive. Returns list of (filename, content)."""
    results = []
    with tempfile.TemporaryDirectory() as tmpdir:
        zip_path = os.path.join(tmpdir, "archive.zip")
        with open(zip_path, "wb") as f:
            f.write(zip_bytes)

        with zipfile.ZipFile(zip_path, "r") as zf:
            zf.extractall(tmpdir)

        for root, dirs, files in os.walk(tmpdir):
            # Skip __MACOSX and hidden directories
            dirs[:] = [d for d in dirs if not d.startswith((".", "__"))]
            for fname in files:
                if fname.startswith("."):
                    continue
                ext = os.path.splitext(fname)[1].lower()
                if ext in SUPPORTED_EXTENSIONS:
                    fpath = os.path.join(root, fname)
                    with open(fpath, "rb") as f:
                        results.append((fname, f.read()))
    return results


def load_corpus() -> None:
    """Load all supported files from corpus, chunk, embed, store in ChromaDB."""
    global _collection

    client = _get_client()

    try:
        client.delete_collection("corpus")
    except Exception:
        pass

    _collection = client.create_collection(
        name="corpus",
        metadata={"hnsw:space": "cosine"},
    )

    model = _get_model()
    all_chunks: list[str] = []
    all_ids: list[str] = []
    all_meta: list[dict] = []

    if not os.path.isdir(CORPUS_DIR):
        os.makedirs(CORPUS_DIR, exist_ok=True)
        return

    for filename in sorted(os.listdir(CORPUS_DIR)):
        filepath = os.path.join(CORPUS_DIR, filename)
        ext = os.path.splitext(filename)[1].lower()
        if ext not in SUPPORTED_EXTENSIONS:
            continue

        text = _read_file(filepath)
        if not text.strip():
            continue

        chunks = _approximate_token_split(text, CHUNK_SIZE, CHUNK_OVERLAP)
        for i, chunk in enumerate(chunks):
            chunk_id = f"{filename}_{i}"
            all_chunks.append(chunk)
            all_ids.append(chunk_id)
            all_meta.append({"source": filename, "chunk_index": i})

    if all_chunks:
        embeddings = model.encode(all_chunks).tolist()
        _collection.add(
            ids=all_ids,
            embeddings=embeddings,
            documents=all_chunks,
            metadatas=all_meta,
        )


def _add_single_file(filename: str, file_bytes: bytes) -> dict:
    """Process a single file: save to corpus and embed."""
    global _collection

    os.makedirs(CORPUS_DIR, exist_ok=True)
    filepath = os.path.join(CORPUS_DIR, filename)

    with open(filepath, "wb") as f:
        f.write(file_bytes)

    text = _read_file(filepath)
    if not text.strip():
        os.remove(filepath)
        return {"filename": filename, "status": "error", "message": "Texte non extractible"}

    chunks = _approximate_token_split(text, CHUNK_SIZE, CHUNK_OVERLAP)
    model = _get_model()

    if _collection is None:
        load_corpus()
        return {"filename": filename, "status": "ok", "chunks": len(chunks)}

    # Remove old chunks from same file if re-uploading
    try:
        existing = _collection.get(where={"source": filename})
        if existing["ids"]:
            _collection.delete(ids=existing["ids"])
    except Exception:
        pass

    chunk_ids = [f"{filename}_{i}" for i in range(len(chunks))]
    metas = [{"source": filename, "chunk_index": i} for i in range(len(chunks))]
    embeddings = model.encode(chunks).tolist()

    _collection.add(
        ids=chunk_ids,
        embeddings=embeddings,
        documents=chunks,
        metadatas=metas,
    )

    return {"filename": filename, "status": "ok", "chunks": len(chunks)}


def add_documents(files: list[tuple[str, bytes]]) -> list[dict]:
    """Add one or more uploaded files. Handles ZIP extraction automatically."""
    results = []
    for filename, file_bytes in files:
        if filename.lower().endswith(".zip"):
            extracted = _extract_zip(file_bytes)
            if not extracted:
                results.append({"filename": filename, "status": "error",
                                "message": "Aucun fichier supporte trouve dans le ZIP"})
                continue
            for inner_name, inner_bytes in extracted:
                results.append(_add_single_file(inner_name, inner_bytes))
        else:
            results.append(_add_single_file(filename, file_bytes))
    return results


def list_documents() -> list[dict]:
    """List all documents in the corpus directory."""
    docs = []
    if not os.path.isdir(CORPUS_DIR):
        return docs
    for filename in sorted(os.listdir(CORPUS_DIR)):
        ext = os.path.splitext(filename)[1].lower()
        if ext in SUPPORTED_EXTENSIONS:
            filepath = os.path.join(CORPUS_DIR, filename)
            size = os.path.getsize(filepath)
            docs.append({"filename": filename, "size": size})
    return docs


def delete_document(filename: str) -> bool:
    """Delete a document from corpus and its embeddings."""
    global _collection
    filepath = os.path.join(CORPUS_DIR, filename)
    if not os.path.isfile(filepath):
        return False

    os.remove(filepath)

    if _collection is not None:
        try:
            existing = _collection.get(where={"source": filename})
            if existing["ids"]:
                _collection.delete(ids=existing["ids"])
        except Exception:
            pass

    return True


def retrieve(query: str, top_k: int = TOP_K) -> list[str]:
    """Retrieve the top_k most relevant chunks for a query."""
    if _collection is None or _collection.count() == 0:
        return []

    model = _get_model()
    query_embedding = model.encode([query]).tolist()
    results = _collection.query(
        query_embeddings=query_embedding,
        n_results=min(top_k, _collection.count()),
    )
    return results["documents"][0] if results["documents"] else []
